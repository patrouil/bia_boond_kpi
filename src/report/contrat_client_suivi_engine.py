import datetime
import logging

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.text.text import TextFrame

from boond.boond_api import BoondApi
from entities.maturite_contrat import MaturiteContrat
from entities.report_definition import ReportDefinition
from query.prod_plan_query import ProductionPlanQuery
from mapper.suivi_client_mapper import SuiviClientMapper
from report.generic_report import GenericReport


class ContratClientSuiviEngine(GenericReport):
    SUIVI_PER_SLIDE = 10

    def __init__(self, api: BoondApi, definition: ReportDefinition):
        super().__init__(api, definition)
        self.logger = logging.getLogger(__name__)
        self.pole_id = definition.pole_id

    def write_suivi_slide(self, suivi_list: [MaturiteContrat], prs: Presentation) -> None:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = 'Suivi de Missions'
        x, y, cx, cy = Cm(1), Cm(2.3), Cm(31), Cm(len(suivi_list) * 0.5)

        shape = slide.shapes.add_table(len(suivi_list) + 1, 4, x, y, cx, cy)
        table = shape.table
        # add title
        table.cell(0, 0).text = 'Date'
        table.cell(0, 1).text = 'Client'
        table.cell(0, 2).text = 'Action'
        table.cell(0, 3).text = 'Reference'

        table.columns[0].width = Cm(3.5)
        table.columns[1].width = Cm(6)
        table.columns[2].width = Cm(6)
        table.columns[3].width = Cm(15.5)

        row = 1
        s = MaturiteContrat()
        for s in suivi_list:
            self.put_text(table.cell(row, 0).text_frame, s.date[:10])
            self.put_text(table.cell(row, 1).text_frame, s.person)
            if (s.qui is not None): self.put_text(table.cell(row, 2).text_frame, s.qui + '\n' + s.typeOf)
            if (s.contenu is not None): self.put_text(table.cell(row, 3).text_frame, s.contenu)
            row = row + 1
        # end for
        return

    def report_pptx(self, prs: Presentation) -> None:
        p = ProductionPlanQuery(self.api, self.pole_id).getProdPlan()

        m = SuiviClientMapper(self.api)
        suivi_list = m.map(p)
        suivi_list = sorted(suivi_list, key=lambda c: c.date, reverse=False)
        row = 0
        while row < len(suivi_list):
            j = min(len(suivi_list), row + self.SUIVI_PER_SLIDE)
            self.write_suivi_slide(suivi_list[row: j], prs)
            row = row + self.SUIVI_PER_SLIDE
        #
        return None
