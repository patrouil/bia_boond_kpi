import datetime
import logging

from boond.boond_api import BoondApi
from entities.maturite_contrat import MaturiteContrat
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.text.text import TextFrame

from entities.report_definition import ReportDefinition
from mapper.contrat_maturite_mapper import MaturiteContratMapper
from query.prod_plan_query import ProductionPlanQuery
from report.generic_report import GenericReport

class ContratMaturityEngine(GenericReport):

    PROJECTS_PER_SLIDE = 10
    def __init__(self, api:BoondApi, definition: ReportDefinition):
        super().__init__(api, definition)
        self.logger = logging.Logger(__name__)
        self.poleId = definition.pole_id

    def _write_matutity_slide(self, kpi_list:[MaturiteContrat], prs:Presentation)->None:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = "Maturité des Contrats"
        x, y, cx, cy = Cm(1), Cm(2.3), Cm(31), Cm(len(kpi_list) * 0.5)

        shape  = slide.shapes.add_table(len(kpi_list)+1, 4, x, y, cx, cy)
        table = shape.table
        # add title
        table.cell(0,0).text = 'Début'
        table.cell(0,1).text = 'Fin'
        table.cell(0,2).text = 'Consultant'
        table.cell(0,3).text = 'Client'

        table.columns[0].width = Cm(3.5)
        table.columns[1].width = Cm(3.5)
        table.columns[2].width = Cm(12)
        table.columns[3].width = Cm(12)

        row = 1
        for kpi in kpi_list:
            self.put_text(table.cell(row, 0).text_frame, kpi.debut)
            self.put_text(table.cell(row, 1).text_frame, kpi.fin)
            self.put_text(table.cell(row, 2).text_frame, kpi.consultant + '\n' + kpi.reference)
            self.put_text(table.cell(row, 3).text_frame, kpi.donneurOrdre + '\n' + kpi.client)
            row = row+1
        #end for
        return
    def report_pptx(self, prs:Presentation) -> None:
        p = ProductionPlanQuery(self.api, self.poleId).getProdPlan()
        m = MaturiteContratMapper(self.api)
        kpi_list = m.map(p)
        kpi_list = sorted(kpi_list, key= lambda c: c.fin)
        row = 0
        while row < len(kpi_list) :
            j = min(len(kpi_list) , row+self.PROJECTS_PER_SLIDE)
            self._write_matutity_slide(kpi_list[row: j], prs)
            row = row+self.PROJECTS_PER_SLIDE
        #
        return None

